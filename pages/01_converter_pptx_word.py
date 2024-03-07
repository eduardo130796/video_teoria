import streamlit as st
from pptx import Presentation
from docx import Document
import io

# Funções auxiliares e a função de conversão
def aplicar_substituicoes(texto, substituicoes):
    for palavra_antiga, palavra_nova in substituicoes:
        texto = texto.replace(palavra_antiga, palavra_nova)
    return texto

def dividir_texto(texto, limite_caracteres, delimitador):
    partes = []
    parte_atual = ""

    for palavra in texto.split():
        if len(parte_atual + ' ' + palavra) > limite_caracteres and parte_atual.endswith(delimitador):
            partes.append(parte_atual)
            parte_atual = palavra
        else:
            if parte_atual:
                parte_atual += ' '
            parte_atual += palavra

    if parte_atual:
        partes.append(parte_atual)

    return partes

def pptx_to_word_with_slide_markers(pptx_memory, word_memory, substituicoes, limite_caracteres, delimitador):
    prs = Presentation(pptx_memory)
    doc = Document()

    for i, slide in enumerate(prs.slides):
        doc.add_paragraph(f"SLIDE: {i+1}")

        if i > 0:
            doc.add_paragraph().add_run().add_break()

        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text
            if title_text:
                title_text = aplicar_substituicoes(title_text, substituicoes)
                doc.add_paragraph("TITLE: " + title_text, style='Heading 1')

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text != title_text:
                shape_text = aplicar_substituicoes(shape.text, substituicoes)
                partes_texto = dividir_texto(shape_text, limite_caracteres, delimitador)
                for parte in partes_texto:
                    doc.add_paragraph(parte)

    doc.save(word_memory)

# Substituições de palavras e funções de apoio
substituicoes = [
    ("Hely", "Elí"),
    ("Di Pietro", "Di piêtro"),
    ("CF", "Constituição Federal"),
    ("nº", "número"),
    ("art.", "artigo"),
    ("Art.", "artigo"),
    ("obs.", "observação"),
    ("Obs.:", "observação"),
    ("J.J.", "José Joaquim"),
    ("j.j.", "José Joaquim"),
    ("habeas", "habias"),
    ("Habeas", "habias"),
    ("corpus", "corpos"),
    ("§", "parágrafo"),
    ("LOA", "Lei Orçamentária Anual"),
]
limite_caracteres = 100  # Limite de caracteres para divisão do texto
delimitador = '.'        # Delimitador para divisão do texto

def streamlit_app():
    st.title("Conversor PPTX para DOCX com Marcadores de Slide")

    pptx_file = st.file_uploader("Escolha o arquivo PPTX", type="pptx")
    if pptx_file is not None:
        pptx_memory = io.BytesIO(pptx_file.getvalue())
        word_memory = io.BytesIO()

        if st.button('Converter PPTX para DOCX'):
            with st.spinner('Convertendo...'):
                # Iniciar a barra de progresso
                progress_bar = st.progress(0)
                for i in range(1, 101):
                    # Atualizando a barra de progresso
                    progress_bar.progress(i)
                pptx_to_word_with_slide_markers(pptx_memory, word_memory, substituicoes, limite_caracteres, delimitador)
                
                # Resetar a barra de progresso após a conclusão
                progress_bar.empty()

            # Exibir mensagem de sucesso
            st.success('Conversão concluída com sucesso!')

            word_memory.seek(0)
            st.download_button(
                label="Baixar arquivo DOCX",
                data=word_memory,
                file_name='documento_convertido.docx',
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

if __name__ == "__main__":
    streamlit_app()
